import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches


def cell_to_indices(cell_ref):
    col_ref, row_ref = ''.join(filter(str.isalpha, cell_ref)), ''.join(filter(str.isdigit, cell_ref))
    col_idx = sum([(ord(char) - ord('A') + 1) * (26 ** idx) for idx, char in enumerate(col_ref[::-1])]) - 1
    row_idx = int(row_ref) - 1
    return row_idx, col_idx


def extract_tables(data):
    tables_references = {
"Overall Traffic Visits": ("A11", "S64"),
        "Overall Traffic Leads": ("A70", "D122"),
        "Overall Traffic Actions": ("A128", "F180"),
        "New Visitors Visits": ("A194", "S247"),
        "New Visitors Leads": ("A253", "D305"),
        "New Visitors Actions": ("A311", "F363"),
        "Return Visitors Visits": ("A377", "S430"),
        "Return Visitors Leads": ("A436", "D488"),
        "Return Visitors Actions": ("A494", "F546"),
        "Desktop Traffic Visits": ("A560", "S613"),
        "Desktop Traffic Leads": ("A619", "D671"),
        "Desktop Traffic Actions": ("A677", "F729"),
        "Mobile Traffic Visits": ("A743", "S796"),
        "Mobile Traffic Leads": ("A802", "D854"),
        "Mobile Traffic Actions": ("A860", "F912"),
        "Tablet Traffic Visits": ("A926", "S979"),
        "Tablet Traffic Leads": ("A985", "D1037"),
        "Tablet Traffic Actions": ("A1043", "F1095"),
        "Paid (SEA) Traffic Visits": ("A1109", "S1162"),
        "Paid (SEA) Traffic Leads": ("A1168", "D1220"),
        "Paid (SEA) Traffic Actions": ("A1226", "F1278"),
        "Organic (SEO) Traffic Visits": ("A1292", "S1345"),
        "Organic (SEO) Traffic Leads": ("A1351", "D1403"),
        "Organic (SEO) Traffic Actions": ("A1409", "F1461"),
        "Social Traffic Visits": ("A1475", "S1527"),
        "Social Traffic Leads": ("A1533", "D1585"),
        "Social Traffic Actions": ("A1591", "F1643"),
        "Direct Traffic Visits": ("A1657", "S1709"),
        "Direct Traffic Leads": ("A1715", "D1767"),
        "Direct Traffic Actions": ("A1773", "F1825"),
        "Session Refresh Traffic Visits": ("A1839", "S1891"),
        "Session Refresh Traffic Leads": ("A1897", "D1949"),
        "Session Refresh Traffic Actions": ("A1955", "F2007"),
        "Refferfal Traffic Visits": ("A2021", "S2073"),
        "Refferfal Traffic Leads": ("A2079", "D2131"),
        "Refferfal Traffic Actions": ("A2137", "F2189"),
    }

    tables_data = {}
    for table_name, (start_cell, end_cell) in tables_references.items():
        start_row, start_col = cell_to_indices(start_cell)
        end_row, end_col = cell_to_indices(end_cell)
        extracted_table = data.iloc[start_row:end_row + 1, start_col:end_col + 1]
        header_row = extracted_table.iloc[0]
        extracted_table = extracted_table[1:]
        extracted_table.columns = header_row

        first_column_name = extracted_table.columns[0]
        extracted_table[first_column_name] = pd.to_datetime(extracted_table[first_column_name], format='%d/%m/%Y')
        extracted_table.rename(columns={first_column_name: 'Week'}, inplace=True)

        tables_data[table_name] = extracted_table
    return tables_data


def create_presentation(dataframes, selected_week):
    prs = Presentation()

    # Extract relevant dataframes based on their names
    leads_df = dataframes['Overall Traffic Leads']
    visits_df = dataframes['Overall Traffic Visits']
    actions_df = dataframes['Overall Traffic Actions']

    # Filter data up to the selected week
    leads_df = leads_df[leads_df['Week'] <= pd.to_datetime(selected_week)]
    visits_df = visits_df[visits_df['Week'] <= pd.to_datetime(selected_week)]
    actions_df = actions_df[actions_df['Week'] <= pd.to_datetime(selected_week)]

    def add_chart_to_slide(df, column, title_text):
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = title_text + " Up To " + selected_week
        chart_data = CategoryChartData()
        chart_data.categories = df['Week'].dt.strftime('%d/%m/%Y')
        chart_data.add_series(title_text, df[column])

        x, y, cx, cy = Inches(1), Inches(1), Inches(8), Inches(6)
        slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        )

    # Add charts to the slide
    add_chart_to_slide(leads_df, leads_df.columns[1], 'Leads')
    add_chart_to_slide(visits_df, visits_df.columns[2], 'Visits')
    add_chart_to_slide(actions_df, actions_df.columns[1], 'Actions')

    # Save the presentation
    prs.save("summary_presentation.pptx")


st.title("Welcome To The CUPRA Co-Pilot")

uploaded_file = st.file_uploader("Choose a CSV file", type="csv")
if uploaded_file is not None:
    data = pd.read_csv(uploaded_file)
    tables_dataframes = extract_tables(data)
    available_dates = list(tables_dataframes["Overall Traffic Visits"]['Week'].dt.strftime('%d/%m/%Y'))
    selected_week = st.selectbox('Select a week for YTD view', available_dates)

    # Generate PowerPoint presentation
    create_presentation(tables_dataframes, selected_week)

    # Provide a download link for the PowerPoint
    st.write("Download the generated [PowerPoint presentation here](summary_presentation.pptx).")

