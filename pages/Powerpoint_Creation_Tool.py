import openai
import streamlit as st
import pandas as pd
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches



def cell_to_indices(cell_ref):
    col_ref, row_ref = ''.join(filter(str.isalpha, cell_ref)), ''.join(filter(str.isdigit, cell_ref))
    col_idx = sum([(ord(char) - ord('A') + 1) * (26 ** idx) for idx, char in enumerate(col_ref[::-1])]) - 1
    row_idx = int(row_ref) - 1
    return row_idx, col_idx
    pass

def extract_tables(data):
    tables_references = {
        "Overall Traffic Visits": ("A11", "S64"),
        "Overall Traffic Leads": ("A70", "D122"),
        "Overall Traffic Actions": ("A128", "F180"),
        "New Visitors Visits": ("A194", "S247"),
        "New Visitors Leads": ("A253", "D305"),
        "New Visitors Actions": ("A311", "F363"),
        "Return Visitors Visits": ("A377", "S430"),
        "Return Visitors Leads": ("A436", "D488"),
        "Return Visitors Actions": ("A494", "F546"),
        "Desktop Traffic Visits": ("A560", "S613"),
        "Desktop Traffic Leads": ("A619", "D671"),
        "Desktop Traffic Actions": ("A677", "F729"),
        "Mobile Traffic Visits": ("A743", "S796"),
        "Mobile Traffic Leads": ("A802", "D854"),
        "Mobile Traffic Actions": ("A860", "F912"),
        "Tablet Traffic Visits": ("A926", "S979"),
        "Tablet Traffic Leads": ("A985", "D1037"),
        "Tablet Traffic Actions": ("A1043", "F1095"),
        "Paid (SEA) Traffic Visits": ("A1109", "S1162"),
        "Paid (SEA) Traffic Leads": ("A1168", "D1220"),
        "Paid (SEA) Traffic Actions": ("A1226", "F1278"),
        "Organic (SEO) Traffic Visits": ("A1292", "S1345"),
        "Organic (SEO) Traffic Leads": ("A1351", "D1403"),
        "Organic (SEO) Traffic Actions": ("A1409", "F1461"),
        "Social Traffic Visits": ("A1475", "S1527"),
        "Social Traffic Leads": ("A1533", "D1585"),
        "Social Traffic Actions": ("A1591", "F1643"),
        "Direct Traffic Visits": ("A1657", "S1709"),
        "Direct Traffic Leads": ("A1715", "D1767"),
        "Direct Traffic Actions": ("A1773", "F1825"),
        "Session Refresh Traffic Visits": ("A1839", "S1891"),
        "Session Refresh Traffic Leads": ("A1897", "D1949"),
        "Session Refresh Traffic Actions": ("A1955", "F2007"),
        "Refferfal Traffic Visits": ("A2021", "S2073"),
        "Refferfal Traffic Leads": ("A2079", "D2131"),
        "Refferfal Traffic Actions": ("A2137", "F2189"),
    }

    tables_data = {}
    for table_name, (start_cell, end_cell) in tables_references.items():
        start_row, start_col = cell_to_indices(start_cell)
        end_row, end_col = cell_to_indices(end_cell)
        extracted_table = data.iloc[start_row:end_row + 1, start_col:end_col + 1]
        extracted_table.rename(columns={extracted_table.columns[0]: 'Week'}, inplace=True)
        extracted_table['Week'] = pd.to_datetime(extracted_table['Week'], format='%d/%m/%Y')
        tables_data[table_name] = extracted_table
    return tables_data

# Defining the create_presentation function in detail

def create_presentation(dataframes, selected_week):
    # Create a new PowerPoint presentation
    prs = Presentation()

    # Extract relevant dataframes
    leads_df = dataframes['TLA Lead Gen.']
    visits_df = dataframes['Overall Traffic Visits']
    actions_df = dataframes['Overall Actions']

    # Filter data up to the selected week
    leads_df = leads_df[leads_df['Week'] <= pd.to_datetime(selected_week)]
    visits_df = visits_df[visits_df['Week'] <= pd.to_datetime(selected_week)]
    actions_df = actions_df[actions_df['Week'] <= pd.to_datetime(selected_week)]

    # Calculate YTD totals
    total_leads_ytd = leads_df['TLA Leads'].sum()
    total_visits_ytd = visits_df['Visits'].sum()
    total_actions_ytd = actions_df['Actions'].sum()

    # Create a new slide with title and content layout
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title = slide.shapes.title
    title.text = "Summary Up To " + selected_week

    # Add YTD totals to the slide
    content = slide.placeholders[1]
    content.text = f"Total Leads YTD: {total_leads_ytd}\nTotal Visits YTD: {total_visits_ytd}\nTotal Actions YTD: {total_actions_ytd}"

    # Create a function to add a chart to the presentation
    def add_chart_to_slide(df, column, title_text, x, y):
        chart_data = CategoryChartData()
        chart_data.categories = df['Week'].dt.strftime('%d/%m/%Y')
        chart_data.add_series(title_text, df[column])

        x, y, cx, cy = Inches(x), Inches(y), Inches(4), Inches(3)
        slide.shapes.add_chart(
            XL_CHART_TYPE.LINE, x, y, cx, cy, chart_data
        )

    # Add charts to the slide
    add_chart_to_slide(leads_df, 'TLA Leads', 'Leads', 0.5, 3)
    add_chart_to_slide(visits_df, 'Visits', 'Visits', 5, 3)
    add_chart_to_slide(actions_df, 'Actions', 'Actions', 0.5, 6)

    # Save the presentation
    prs.save("/mnt/data/summary_presentation.pptx")

# Placeholder output to confirm the function definition
"create_presentation function defined successfully!"


st.title("Welcome To The CUPRA Co-Pilot")

uploaded_file = st.file_uploader("Choose a CSV file", type="csv")
if uploaded_file is not None:
    data = pd.read_csv(uploaded_file)
    tables_dataframes = extract_tables(data)
    available_dates = list(tables_dataframes["Overall Traffic Visits"]['Week'].dt.strftime('%d/%m/%Y'))
    selected_week = st.selectbox('Select a week for YTD view', available_dates)

    # Generate PowerPoint presentation
    create_presentation(tables_dataframes, selected_week)

    # Provide a download link for the PowerPoint
    st.write("Download the generated [PowerPoint presentation here](/mnt/data/summary_presentation.pptx).")

